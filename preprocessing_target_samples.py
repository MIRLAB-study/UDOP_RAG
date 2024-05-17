import os
import argparse
import json
from tqdm import tqdm

import fitz
from PIL import Image
from pptx import Presentation # mdj: pip install python-pptx==0.6.23

def parse_pdf(pdf_file="./target_samples/240314_문서분석_산학.pdf"):
    # Split the base name and extension
    output_directory_path, _ = os.path.splitext(pdf_file)

    if not os.path.exists(output_directory_path):
        os.makedirs(output_directory_path)
    
    # Open the PDF file
    pdf_document = fitz.open(pdf_file)

    # Iterate through each page and convert to an image
    abs_output_path = os.path.abspath(output_directory_path)
    for page_number in tqdm(range(pdf_document.page_count)):
        # Get the page
        page = pdf_document[page_number]

        # Convert the page to an image
        pix = page.get_pixmap()

        # Create a Pillow Image object from the pixmap
        image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

        # Save the image
        image.save(f"{abs_output_path}/page_{page_number + 1}.png")

    # Close the PDF file
    pdf_document.close()

    image_paths = []
    for img_path in os.listdir(output_directory_path):
        image_paths.append(str(os.path.join(output_directory_path, img_path)))
    return output_directory_path, image_paths

def parse_pptx(pptx_file="./target_samples/240314_문서분석_산학.pptx"):
    # from IPython import embed; embed()

    # Split the base name and extension
    output_directory_path, _ = os.path.splitext(pptx_file)

    # JSON file name
    json_path = os.path.join(output_directory_path, 'Bbox_metadata.json')

    # Load the presentation
    prs = Presentation(pptx_file)

    # Dictionary to hold the coordinates for each slide
    slide_Bbox_metadata = {}

    # Iterate over each slide in the presentation
    for i, slide in tqdm(enumerate(prs.slides)):
        if i<26: continue# mdj: Starting point of category slide in Hyundai PPT

        # Reset counting
        Bbox_name_list = {
            'title': 0, 
            'text': 0, 
            'table': 0, 
            'image': 0, 
            'allow': 0, 
            'milestone': 0, 
            'figure': 0,
            '': 0, # Unnamed Bbox
            'flowchart': 0,
            'attachment': 0,
            'graph': 0,
            'legend': 0,
            'number': 0,
            'Q. 첨부파일 열어 내용 확인 후 학습 가능한가?': 0,
        }
        
        # Dictionary to hold the coordinates of text boxes for the current slide
        textbox_coordinates = {}

        # Iterate over each shape in the slide
        for shape in slide.shapes:
            # Check if the shape has text (i.e., it is a text box)
            if shape.has_text_frame:
                # Extract coordinates and dimensions
                x, y = shape.left, shape.top
                width, height = shape.width, shape.height

                # Extract text
                text = shape.text
                
                # Store the coordinates in a dictionary with text as the key
                textbox_coordinates[f"{text}_{Bbox_name_list[text]}"] = (x, y, width, height)

                # Counting text in each slide
                Bbox_name_list[text] += 1

        # Add the current slide's text box coordinates to the main dictionary
        slide_key = f"slide_{i + 1}"
        slide_Bbox_metadata[slide_key] = textbox_coordinates
    
    # Save the Bbox metadata to a JSON file
    with open(json_path, 'w', encoding='utf-8') as json_file:
        json.dump(slide_Bbox_metadata, json_file, ensure_ascii=False, indent=4)

def main(args):
    if args.pdf_to_image:
        parse_pdf(args.pdf_path)
    elif args.save_metadata_from_pptx:
        parse_pptx(args.pptx_path)
    else:
        raise NotImplementedError()

if __name__ == '__main__':
    # 1. create parser
    parser = argparse.ArgumentParser(
        description='This Python program preprocesses the target data provided by Hyundai.'
    )

    # 2. add arguments to parser
    parser.add_argument(
        '--pdf-to-image', 
        action="store_true", 
        help='Convert and save each page of a PDF file as an image file.'
    )
    parser.add_argument(
        '--save-metadata-from-pptx', 
        action="store_true", 
        help='''
            Extract Bbox metadata ('x', 'y', 'width', 'height') from each slide in a PPT file and save it as a JSON file.
            '''
    )
    parser.add_argument(
        '--pdf-path', 
        type=str, 
        default='./target_samples/240314_문서분석_산학.pdf', 
        help='Path to PDF file'
    )
    parser.add_argument(
        '--pptx-path', 
        type=str, 
        default='./target_samples/240314_문서분석_산학.pptx', 
        help='Path to PPT file'
    )

    # 3. parse arguments
    args = parser.parse_args()

    # 4. use arguments
    print (args)
    assert args.pdf_to_image or args.save_metadata_from_pptx, 'Choose pdf-to-image or save-metadata-from-pptx'
    main(args)
    # python preprocessing_target_samples.py --pdf-to-image --pdf-path './target_samples/240314_문서분석_산학.pdf'
    # python preprocessing_target_samples.py --save-metadata-from-pptx --pptx-path './target_samples/240314_문서분석_산학.pptx'